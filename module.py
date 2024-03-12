# dependency
# nltk 3.7
# konlpy 0.6.0
# wordcloud 1.8.2.2
# pptx 0.6.21
# pip install nltk
# pip install konlpy
# pip install wordcloud
# pip install python-pptx
# jre 설치 필요
# command code : 쳐야하는 코드
# python simple_keywording "[파이썬 파일과 같은 위치에 있는 대상 ppt이름]" <- 쌍따옴표 2개는 중요!
import sys
import os
import re
import string
import nltk
import time
import pandas
import numpy
import logging
import re
import threading
import multiprocessing
from pptx import Presentation
from konlpy.tag import Okt
from konlpy.tag import Komoran, Hannanum
from collections import Counter
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from wordcloud import WordCloud
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError, InvalidXmlError, PythonPptxError

preprocessing_module_path = os.path.abspath(os.path.dirname(__file__))
stopwords_path = os.path.join(preprocessing_module_path, "module_stopwords.txt")


logging.basicConfig(filename="./extracting.log", level=logging.ERROR)

# load komoran
komoran = Komoran()

# stopword embedding
# stopwords ( text from : https://raw.githubusercontent.com/yoonkt200/FastCampusDataset/master/korean_stopwords.txt )
with open(stopwords_path, "r", encoding="utf-8") as f:
    list_file = f.readlines()
stopwords = list(map(lambda x: x.replace("\n", ""), list_file))

# function : extract string from pptx
# input : string = ppt_name where file extension is .pptx
# output : string = low ppt text
# dependancy : pptx


def extract_pptx_cell_text(shape):
    ret_cell_text = ""
    row_count = len(shape.table.rows)
    col_count = len(shape.table.columns)
    for _r in range(0, row_count):
        for _c in range(0, col_count):
            ret_cell_text += shape.table.cell(_r, _c).text + " "
    return ret_cell_text


def extract_pptx_low_text_str(ppt_name):
    low_text_list = ""
    input = Presentation(ppt_name)
    for slide in input.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                low_text_list += extract_pptx_cell_text(shape)
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                low_text_list += paragraph.text + " "
    return low_text_list.rstrip()


def extract_ppt_low_text_str_threading(ppt_name, start, end):
    low_text_list = ""
    print(f"Current thread ID: {threading.get_ident()}")
    scope_slides = list(Presentation(ppt_name).slides)[start:end+1]
    for slide in scope_slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                low_text_list += extract_pptx_cell_text(shape)
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                # print(paragraph.text)
                low_text_list += paragraph.text + " "
    return low_text_list.rstrip()


# basic preprocessing
# input : String
# output : String
# dependancy : re


def preprocess(text):
    text = text.strip()
    text = re.compile("[%s]" % re.escape(string.punctuation)).sub(" ", text)
    text = re.sub("\s+", " ", text)
    text = re.sub(
        r"\b\d+\b|[^\w\s]",
        lambda match: match.group() if match.group().startswith("\b") else "",
        text,
    )
    text = re.sub(r"[^\w\s]", " ", str(text).strip())
    text = re.sub(r"\s+", " ", text)
    return text


# function : extracting keyword from komoran
# input : String
# output : String
# dependancy : konlpy,


def final(text):
    n = []
    word = komoran.nouns(text)
    p = komoran.pos(text)
    # hannanum_word = Hannanum().nouns(text)
    # print(hannanum_word)
    # word +=hannanum_word
    # hannanum_p = Hannanum().pos(text)
    # print(hannanum_p)
    for pos in p:
        if pos[1] in ["SL"]:
            word.append(pos[0])
    # for hpos in hannanum_p:
    #     if hpos[1] in ["SL"]:
    #         word.append(hpos[0])
    for w in word:
        if len(w) > 1 and w not in stopwords:
            n.append(w)
    # print(word)
    return " ".join(n)


# function : all preprocessing file
# input : string = pptx title
# output : string = keyword with " "


def final_preprocessing(ppt_name):
    return final(preprocess(extract_pptx_low_text_str(ppt_name)))


# link : final_preprocessing
# function : extract preprocessing pptx text data
# input : String = file_path : any path where [absolute_path] or [relative_path]
# output : dict = preprocessed keyword


def file2data(file_path):
    name_with_keyword = {}
    name_with_keyword[file_path] = final_preprocessing(file_path)
    return name_with_keyword


def errorFileDetection(file_path):
    try:
        Presentation(file_path)
        logging.info(f"extrating right file in {file_path}")
        return True
    except PackageNotFoundError as error:
        logging.exception(f"extracting file Exception in {file_path}")
        return False
    except InvalidXmlError as error:
        logging.exception(f"extracting file Exception in {file_path}")
        return False
    except PythonPptxError as error:
        logging.exception(f"extracting file Exception in {file_path}")
        return False
    except Exception as error:
        logging.exception(f"other Exception in {file_path}")
        return False


# 대분류
# link : file2data
# function : extract all folder using recursive method pptx text data
# input : String = any_folder_path
# output : dict = {file_path : keyword_data}


def folder2data_threading(any_folder_path):
    return_dic = {}
    f_list = os.listdir(any_folder_path)
    pattern = r"\.pptx$"
    print(
        "   now extracting in {} folder wiht length {} loading...".format(
            any_folder_path, len(f_list)
        )
    )
    start = time.time()
    for file_path in f_list:
        absoulte_folder_path = os.path.join(any_folder_path, file_path)
        if (
            errorFileDetection(absoulte_folder_path)
            and re.search(pattern, file_path)
            and file_path.find("~$") == -1
            and os.stat(os.path.join(any_folder_path, file_path)).st_size > 0
        ):
            print(" this is in Threading Part")
            parent_path = os.path.join(any_folder_path, file_path)
            # Thread Programming
            print("     out Thread")
            if len(Presentation(parent_path).slides) > 100:
                print("     in Thread")
                cpu_count = multiprocessing.cpu_count()
                print(cpu_count)
                thread_epoch = len(Presentation(parent_path).slides) // cpu_count
                print(thread_epoch)
                thread_result = []
                threads = []
                thread_result_str = ""
                for epoch in range(cpu_count):
                    t = threading.Thread(
                        target=lambda: thread_result.append(
                            extract_ppt_low_text_str_threading(
                                parent_path,
                                epoch * thread_epoch,
                                (epoch + 1) * thread_epoch,
                            )
                        )
                    )
                    threads.append(t)
                    t.start()
                for runned_thread in threads:
                    runned_thread.join()
                for result in thread_result:
                    thread_result_str += result + " "
                temp_dic = {}
                temp_dic[parent_path] = thread_result_str
                return_dic.update(temp_dic)

            else:
                parent_path = os.path.join(any_folder_path, file_path)
                return_dic.update(file2data(parent_path))
        elif os.path.isdir(os.path.join(any_folder_path, file_path)):
            child_path = os.path.join(any_folder_path, file_path)
            return_dic.update(folder2data_threading(child_path))
    print("     end extracting in {} folder...".format(any_folder_path))
    end = time.time()
    print("     extracting folder time {} sec".format(end - start))
    return return_dic





def folder2data(any_folder_path):
    return_dic = {}
    f_list = os.listdir(any_folder_path)
    pattern = r"\.pptx$"
    print(
        "   now extracting in {} folder wiht length {} loading...".format(
            any_folder_path, len(f_list)
        )
    )
    start = time.time()
    for file_path in f_list:
        absoulte_folder_path = os.path.join(any_folder_path, file_path)
        if os.path.isdir(os.path.join(any_folder_path, file_path)):
            child_path = os.path.join(any_folder_path, file_path)
            return_dic.update(folder2data(child_path))
        elif (
            errorFileDetection(absoulte_folder_path)
            and re.search(pattern, file_path)
            and file_path.find("~$") == -1
            and os.stat(os.path.join(any_folder_path, file_path)).st_size > 0
        ):
            parent_path = os.path.join(any_folder_path, file_path)
            return_dic.update(file2data(parent_path))
    print("     end extracting in {} folder...".format(any_folder_path))
    end = time.time()
    print("     extracting folder time {} sec".format(end - start))
    return return_dic
